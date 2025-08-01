"""
Advanced RAG (Retrieval-Augmented Generation) System
Enterprise-grade document processing and retrieval with vector database
"""

import os
import json
import hashlib
import asyncio
from typing import List, Dict, Any, Optional, Tuple, Union
from dataclasses import dataclass, field
from datetime import datetime
from pathlib import Path
import numpy as np
from concurrent.futures import ThreadPoolExecutor, ProcessPoolExecutor
import pickle
import sqlite3
import faiss
import torch
import torch.nn.functional as F
from transformers import (
    AutoTokenizer, 
    AutoModel,
    pipeline,
    DPRQuestionEncoder,
    DPRContextEncoder,
    RagTokenizer,
    RagRetriever,
    RagSequenceForGeneration
)
from sentence_transformers import SentenceTransformer
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.document_loaders import (
    PyPDFLoader,
    UnstructuredWordDocumentLoader,
    UnstructuredPowerPointLoader,
    UnstructuredExcelLoader,
    TextLoader,
    JSONLoader,
    CSVLoader,
    UnstructuredMarkdownLoader,
    UnstructuredHTMLLoader
)
from PyQt6.QtCore import QObject, pyqtSignal, QThread, QMutex, QMutexLocker
import chromadb
from chromadb.config import Settings
import weaviate
from qdrant_client import QdrantClient
from qdrant_client.models import Distance, VectorParams, PointStruct

# Document processing
import fitz  # PyMuPDF
import docx
import openpyxl
import pptx
from bs4 import BeautifulSoup
import markdown
import pytesseract
from PIL import Image
import speech_recognition as sr
from pydub import AudioSegment

# NLP and preprocessing
import spacy
import nltk
from nltk.corpus import stopwords
from nltk.tokenize import word_tokenize
from nltk.stem import WordNetLemmatizer

# Download required NLTK data
nltk.download('punkt', quiet=True)
nltk.download('stopwords', quiet=True)
nltk.download('wordnet', quiet=True)

# Load spaCy model
try:
    nlp = spacy.load("en_core_web_sm")
except:
    os.system("python -m spacy download en_core_web_sm")
    nlp = spacy.load("en_core_web_sm")


@dataclass
class Document:
    """Enterprise document structure"""
    id: str
    content: str
    metadata: Dict[str, Any] = field(default_factory=dict)
    embeddings: Optional[np.ndarray] = None
    chunks: List['DocumentChunk'] = field(default_factory=list)
    processed_at: datetime = field(default_factory=datetime.now)
    source: str = ""
    doc_type: str = ""
    language: str = "en"
    permissions: Dict[str, Any] = field(default_factory=dict)
    version: int = 1
    hash: str = ""
    
    def __post_init__(self):
        if not self.id:
            self.id = hashlib.sha256(self.content.encode()).hexdigest()[:16]
        if not self.hash:
            self.hash = hashlib.sha256(self.content.encode()).hexdigest()


@dataclass
class DocumentChunk:
    """Document chunk for processing"""
    id: str
    document_id: str
    content: str
    embedding: Optional[np.ndarray] = None
    metadata: Dict[str, Any] = field(default_factory=dict)
    start_idx: int = 0
    end_idx: int = 0
    chunk_idx: int = 0


@dataclass
class RetrievalResult:
    """Search result structure"""
    document_id: str
    chunk_id: str
    content: str
    score: float
    metadata: Dict[str, Any] = field(default_factory=dict)
    highlights: List[Tuple[int, int]] = field(default_factory=list)


class DocumentProcessor(QThread):
    """Advanced document processing engine"""
    
    progress = pyqtSignal(int)
    status = pyqtSignal(str)
    document_processed = pyqtSignal(Document)
    error_occurred = pyqtSignal(str)
    
    def __init__(self):
        super().__init__()
        self.supported_formats = {
            '.pdf': self._process_pdf,
            '.docx': self._process_docx,
            '.doc': self._process_doc,
            '.pptx': self._process_pptx,
            '.xlsx': self._process_xlsx,
            '.txt': self._process_text,
            '.md': self._process_markdown,
            '.html': self._process_html,
            '.json': self._process_json,
            '.csv': self._process_csv,
            '.png': self._process_image,
            '.jpg': self._process_image,
            '.jpeg': self._process_image,
            '.mp3': self._process_audio,
            '.wav': self._process_audio,
            '.m4a': self._process_audio
        }
        
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            separators=["\n\n", "\n", " ", ""]
        )
        
        self.lemmatizer = WordNetLemmatizer()
        self.stop_words = set(stopwords.words('english'))
        
    def process_document(self, file_path: str, metadata: Dict[str, Any] = None) -> Document:
        """Process document with enterprise features"""
        self.status.emit(f"Processing {file_path}")
        
        file_path = Path(file_path)
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
            
        # Detect file type
        ext = file_path.suffix.lower()
        if ext not in self.supported_formats:
            raise ValueError(f"Unsupported file format: {ext}")
            
        # Extract content
        content = self.supported_formats[ext](file_path)
        
        # Create document
        doc = Document(
            id=hashlib.sha256(str(file_path).encode()).hexdigest()[:16],
            content=content,
            source=str(file_path),
            doc_type=ext[1:],
            metadata=metadata or {}
        )
        
        # Advanced NLP processing
        doc = self._enhance_with_nlp(doc)
        
        # Create chunks
        chunks = self._create_chunks(doc)
        doc.chunks = chunks
        
        self.document_processed.emit(doc)
        return doc
        
    def _process_pdf(self, file_path: Path) -> str:
        """Extract text from PDF with OCR fallback"""
        text = ""
        try:
            pdf = fitz.open(str(file_path))
            for page_num, page in enumerate(pdf):
                page_text = page.get_text()
                if not page_text.strip():
                    # OCR fallback
                    pix = page.get_pixmap()
                    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    page_text = pytesseract.image_to_string(img)
                text += f"\n[Page {page_num + 1}]\n{page_text}"
                self.progress.emit(int((page_num + 1) / len(pdf) * 100))
            pdf.close()
        except Exception as e:
            self.error_occurred.emit(f"PDF processing error: {str(e)}")
            # Fallback to langchain
            loader = PyPDFLoader(str(file_path))
            pages = loader.load()
            text = "\n".join([p.page_content for p in pages])
        return text
        
    def _process_docx(self, file_path: Path) -> str:
        """Extract text from DOCX"""
        doc = docx.Document(str(file_path))
        text = []
        for para in doc.paragraphs:
            text.append(para.text)
        # Extract from tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    text.append(cell.text)
        return "\n".join(text)
        
    def _process_doc(self, file_path: Path) -> str:
        """Extract text from DOC"""
        loader = UnstructuredWordDocumentLoader(str(file_path))
        docs = loader.load()
        return "\n".join([d.page_content for d in docs])
        
    def _process_pptx(self, file_path: Path) -> str:
        """Extract text from PowerPoint"""
        prs = pptx.Presentation(str(file_path))
        text = []
        for slide_num, slide in enumerate(prs.slides):
            text.append(f"\n[Slide {slide_num + 1}]\n")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
        
    def _process_xlsx(self, file_path: Path) -> str:
        """Extract text from Excel"""
        wb = openpyxl.load_workbook(str(file_path), data_only=True)
        text = []
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            text.append(f"\n[Sheet: {sheet_name}]\n")
            for row in sheet.iter_rows(values_only=True):
                row_text = "\t".join([str(cell) if cell else "" for cell in row])
                if row_text.strip():
                    text.append(row_text)
        return "\n".join(text)
        
    def _process_text(self, file_path: Path) -> str:
        """Process plain text"""
        return file_path.read_text(encoding='utf-8', errors='ignore')
        
    def _process_markdown(self, file_path: Path) -> str:
        """Process markdown"""
        md_text = file_path.read_text(encoding='utf-8')
        html = markdown.markdown(md_text)
        soup = BeautifulSoup(html, 'html.parser')
        return soup.get_text()
        
    def _process_html(self, file_path: Path) -> str:
        """Process HTML"""
        html = file_path.read_text(encoding='utf-8', errors='ignore')
        soup = BeautifulSoup(html, 'html.parser')
        # Remove script and style elements
        for script in soup(["script", "style"]):
            script.decompose()
        return soup.get_text()
        
    def _process_json(self, file_path: Path) -> str:
        """Process JSON"""
        data = json.loads(file_path.read_text())
        return json.dumps(data, indent=2)
        
    def _process_csv(self, file_path: Path) -> str:
        """Process CSV"""
        loader = CSVLoader(str(file_path))
        docs = loader.load()
        return "\n".join([d.page_content for d in docs])
        
    def _process_image(self, file_path: Path) -> str:
        """OCR for images"""
        img = Image.open(file_path)
        text = pytesseract.image_to_string(img)
        return f"[Image: {file_path.name}]\n{text}"
        
    def _process_audio(self, file_path: Path) -> str:
        """Speech to text for audio"""
        r = sr.Recognizer()
        # Convert to WAV if needed
        audio = AudioSegment.from_file(str(file_path))
        wav_path = file_path.with_suffix('.wav')
        audio.export(str(wav_path), format="wav")
        
        with sr.AudioFile(str(wav_path)) as source:
            audio_data = r.record(source)
            try:
                text = r.recognize_google(audio_data)
            except:
                text = "[Audio transcription failed]"
                
        if wav_path != file_path:
            wav_path.unlink()
            
        return f"[Audio: {file_path.name}]\n{text}"
        
    def _enhance_with_nlp(self, doc: Document) -> Document:
        """Advanced NLP enhancement"""
        # Named entity recognition
        doc_nlp = nlp(doc.content[:1000000])  # Limit for performance
        
        entities = [(ent.text, ent.label_) for ent in doc_nlp.ents]
        doc.metadata['entities'] = entities
        
        # Key phrases extraction
        keywords = self._extract_keywords(doc.content)
        doc.metadata['keywords'] = keywords
        
        # Language detection
        doc.language = self._detect_language(doc.content)
        
        # Summary generation (if content is long)
        if len(doc.content) > 5000:
            doc.metadata['summary'] = self._generate_summary(doc.content)
            
        return doc
        
    def _extract_keywords(self, text: str, top_k: int = 10) -> List[str]:
        """Extract keywords using TF-IDF"""
        from sklearn.feature_extraction.text import TfidfVectorizer
        
        # Preprocess
        tokens = word_tokenize(text.lower())
        tokens = [self.lemmatizer.lemmatize(t) for t in tokens 
                 if t.isalnum() and t not in self.stop_words and len(t) > 2]
        
        if not tokens:
            return []
            
        # TF-IDF
        processed_text = ' '.join(tokens)
        vectorizer = TfidfVectorizer(max_features=top_k)
        try:
            vectorizer.fit_transform([processed_text])
            keywords = vectorizer.get_feature_names_out()
            return list(keywords)
        except:
            return tokens[:top_k]
            
    def _detect_language(self, text: str) -> str:
        """Detect document language"""
        try:
            from langdetect import detect
            return detect(text[:1000])
        except:
            return "en"
            
    def _generate_summary(self, text: str) -> str:
        """Generate document summary"""
        try:
            summarizer = pipeline("summarization", model="facebook/bart-large-cnn")
            summary = summarizer(text[:1024], max_length=150, min_length=50, do_sample=False)
            return summary[0]['summary_text']
        except:
            # Fallback to extractive summary
            sentences = text.split('.')[:5]
            return '. '.join(sentences)
            
    def _create_chunks(self, doc: Document) -> List[DocumentChunk]:
        """Create semantic chunks"""
        chunks = []
        text_chunks = self.text_splitter.split_text(doc.content)
        
        for idx, chunk_text in enumerate(text_chunks):
            chunk = DocumentChunk(
                id=f"{doc.id}_chunk_{idx}",
                document_id=doc.id,
                content=chunk_text,
                chunk_idx=idx,
                metadata={
                    'source': doc.source,
                    'doc_type': doc.doc_type,
                    'chunk_size': len(chunk_text)
                }
            )
            chunks.append(chunk)
            
        return chunks


class VectorDatabase:
    """Enterprise vector database with multiple backend support"""
    
    def __init__(self, backend: str = "faiss", **kwargs):
        self.backend = backend
        self.index = None
        self.metadata_store = {}
        self.dimension = kwargs.get('dimension', 768)
        
        if backend == "faiss":
            self._init_faiss(**kwargs)
        elif backend == "chromadb":
            self._init_chromadb(**kwargs)
        elif backend == "weaviate":
            self._init_weaviate(**kwargs)
        elif backend == "qdrant":
            self._init_qdrant(**kwargs)
        else:
            raise ValueError(f"Unsupported backend: {backend}")
            
    def _init_faiss(self, **kwargs):
        """Initialize FAISS index"""
        index_type = kwargs.get('index_type', 'IVF')
        
        if index_type == 'Flat':
            self.index = faiss.IndexFlatL2(self.dimension)
        elif index_type == 'IVF':
            nlist = kwargs.get('nlist', 100)
            quantizer = faiss.IndexFlatL2(self.dimension)
            self.index = faiss.IndexIVFFlat(quantizer, self.dimension, nlist)
        elif index_type == 'HNSW':
            M = kwargs.get('M', 32)
            self.index = faiss.IndexHNSWFlat(self.dimension, M)
        else:
            self.index = faiss.IndexFlatL2(self.dimension)
            
        # Enable GPU if available
        if kwargs.get('use_gpu', True) and torch.cuda.is_available():
            res = faiss.StandardGpuResources()
            self.index = faiss.index_cpu_to_gpu(res, 0, self.index)
            
    def _init_chromadb(self, **kwargs):
        """Initialize ChromaDB"""
        self.client = chromadb.Client(Settings(
            chroma_db_impl="duckdb+parquet",
            persist_directory=kwargs.get('persist_dir', "./chromadb")
        ))
        self.collection = self.client.create_collection(
            name=kwargs.get('collection_name', 'documents'),
            metadata={"hnsw:space": "cosine"}
        )
        
    def _init_weaviate(self, **kwargs):
        """Initialize Weaviate"""
        self.client = weaviate.Client(
            url=kwargs.get('url', 'http://localhost:8080'),
            auth_client_secret=weaviate.AuthApiKey(api_key=kwargs.get('api_key', ''))
        )
        
    def _init_qdrant(self, **kwargs):
        """Initialize Qdrant"""
        self.client = QdrantClient(
            host=kwargs.get('host', 'localhost'),
            port=kwargs.get('port', 6333)
        )
        self.collection_name = kwargs.get('collection_name', 'documents')
        
        # Create collection
        self.client.recreate_collection(
            collection_name=self.collection_name,
            vectors_config=VectorParams(size=self.dimension, distance=Distance.COSINE)
        )
        
    def add_vectors(self, vectors: np.ndarray, metadata: List[Dict[str, Any]]):
        """Add vectors to database"""
        if self.backend == "faiss":
            if hasattr(self.index, 'train') and not self.index.is_trained:
                self.index.train(vectors)
            start_idx = self.index.ntotal
            self.index.add(vectors)
            for i, meta in enumerate(metadata):
                self.metadata_store[start_idx + i] = meta
                
        elif self.backend == "chromadb":
            ids = [meta.get('id', str(i)) for i, meta in enumerate(metadata)]
            self.collection.add(
                embeddings=vectors.tolist(),
                metadatas=metadata,
                ids=ids
            )
            
        elif self.backend == "qdrant":
            points = [
                PointStruct(
                    id=i,
                    vector=vector.tolist(),
                    payload=meta
                )
                for i, (vector, meta) in enumerate(zip(vectors, metadata))
            ]
            self.client.upsert(
                collection_name=self.collection_name,
                points=points
            )
            
    def search(self, query_vector: np.ndarray, k: int = 10, filter: Dict[str, Any] = None) -> List[Tuple[int, float, Dict]]:
        """Search similar vectors"""
        if self.backend == "faiss":
            D, I = self.index.search(query_vector.reshape(1, -1), k)
            results = []
            for i, (idx, dist) in enumerate(zip(I[0], D[0])):
                if idx != -1:
                    meta = self.metadata_store.get(idx, {})
                    if not filter or all(meta.get(k) == v for k, v in filter.items()):
                        results.append((idx, float(dist), meta))
            return results
            
        elif self.backend == "chromadb":
            results = self.collection.query(
                query_embeddings=[query_vector.tolist()],
                n_results=k,
                where=filter
            )
            return [(i, d, m) for i, d, m in zip(
                range(len(results['ids'][0])),
                results['distances'][0],
                results['metadatas'][0]
            )]
            
        elif self.backend == "qdrant":
            search_result = self.client.search(
                collection_name=self.collection_name,
                query_vector=query_vector.tolist(),
                limit=k,
                query_filter=filter
            )
            return [(r.id, r.score, r.payload) for r in search_result]
            
    def save(self, path: str):
        """Save index to disk"""
        if self.backend == "faiss":
            faiss.write_index(self.index, f"{path}.faiss")
            with open(f"{path}.meta", 'wb') as f:
                pickle.dump(self.metadata_store, f)
                
    def load(self, path: str):
        """Load index from disk"""
        if self.backend == "faiss":
            self.index = faiss.read_index(f"{path}.faiss")
            with open(f"{path}.meta", 'rb') as f:
                self.metadata_store = pickle.load(f)


class EmbeddingEngine(QObject):
    """Multi-model embedding engine"""
    
    embedding_complete = pyqtSignal(str, np.ndarray)
    batch_complete = pyqtSignal(list)
    error_occurred = pyqtSignal(str)
    
    def __init__(self, model_name: str = "sentence-transformers/all-mpnet-base-v2"):
        super().__init__()
        self.model_name = model_name
        self.device = torch.device('cuda' if torch.cuda.is_available() else 'cpu')
        
        # Load models
        if "sentence-transformers" in model_name:
            self.model = SentenceTransformer(model_name)
            self.model.to(self.device)
        else:
            self.tokenizer = AutoTokenizer.from_pretrained(model_name)
            self.model = AutoModel.from_pretrained(model_name).to(self.device)
            
        # Optimization
        if torch.cuda.is_available():
            self.model = torch.compile(self.model)
            
    def embed_text(self, text: str) -> np.ndarray:
        """Generate embedding for text"""
        if hasattr(self.model, 'encode'):
            # Sentence transformers
            embedding = self.model.encode(text, convert_to_numpy=True)
        else:
            # Hugging Face transformers
            inputs = self.tokenizer(text, return_tensors='pt', truncation=True, 
                                  padding=True, max_length=512).to(self.device)
            with torch.no_grad():
                outputs = self.model(**inputs)
                # Mean pooling
                embeddings = outputs.last_hidden_state
                mask = inputs['attention_mask'].unsqueeze(-1).expand(embeddings.size()).float()
                masked_embeddings = embeddings * mask
                summed = torch.sum(masked_embeddings, 1)
                summed_mask = torch.clamp(mask.sum(1), min=1e-9)
                embedding = (summed / summed_mask).cpu().numpy()[0]
                
        self.embedding_complete.emit(text[:50], embedding)
        return embedding
        
    def embed_batch(self, texts: List[str], batch_size: int = 32) -> np.ndarray:
        """Batch embedding generation"""
        embeddings = []
        
        for i in range(0, len(texts), batch_size):
            batch = texts[i:i+batch_size]
            
            if hasattr(self.model, 'encode'):
                batch_embeddings = self.model.encode(batch, convert_to_numpy=True, 
                                                   show_progress_bar=False)
            else:
                # Manual batching for HF models
                batch_embeddings = []
                for text in batch:
                    emb = self.embed_text(text)
                    batch_embeddings.append(emb)
                batch_embeddings = np.array(batch_embeddings)
                
            embeddings.extend(batch_embeddings)
            
        embeddings = np.array(embeddings)
        self.batch_complete.emit(embeddings.tolist())
        return embeddings


class RAGPipeline(QObject):
    """Complete RAG pipeline with enterprise features"""
    
    retrieval_complete = pyqtSignal(list)
    generation_complete = pyqtSignal(str)
    pipeline_complete = pyqtSignal(dict)
    error_occurred = pyqtSignal(str)
    
    def __init__(self, vector_db: VectorDatabase, embedding_engine: EmbeddingEngine):
        super().__init__()
        self.vector_db = vector_db
        self.embedding_engine = embedding_engine
        self.device = torch.device('cuda' if torch.cuda.is_available() else 'cpu')
        
        # Load RAG model
        self.tokenizer = RagTokenizer.from_pretrained("facebook/rag-token-nq")
        self.retriever = RagRetriever.from_pretrained("facebook/rag-token-nq")
        self.model = RagSequenceForGeneration.from_pretrained("facebook/rag-token-nq").to(self.device)
        
        # Document processor
        self.doc_processor = DocumentProcessor()
        
        # Cache
        self.cache = {}
        self.cache_size = 1000
        
    def add_documents(self, file_paths: List[str], metadata: List[Dict[str, Any]] = None):
        """Add documents to RAG system"""
        documents = []
        
        for i, file_path in enumerate(file_paths):
            meta = metadata[i] if metadata else {}
            doc = self.doc_processor.process_document(file_path, meta)
            documents.append(doc)
            
        # Generate embeddings for all chunks
        all_chunks = []
        all_metadata = []
        
        for doc in documents:
            for chunk in doc.chunks:
                all_chunks.append(chunk.content)
                all_metadata.append({
                    'doc_id': doc.id,
                    'chunk_id': chunk.id,
                    'source': doc.source,
                    'doc_type': doc.doc_type,
                    **chunk.metadata
                })
                
        # Batch embed
        embeddings = self.embedding_engine.embed_batch(all_chunks)
        
        # Add to vector database
        self.vector_db.add_vectors(embeddings, all_metadata)
        
        return documents
        
    def retrieve(self, query: str, k: int = 5, filters: Dict[str, Any] = None) -> List[RetrievalResult]:
        """Retrieve relevant documents"""
        # Check cache
        cache_key = f"{query}_{k}_{str(filters)}"
        if cache_key in self.cache:
            return self.cache[cache_key]
            
        # Generate query embedding
        query_embedding = self.embedding_engine.embed_text(query)
        
        # Search vector database
        results = self.vector_db.search(query_embedding, k=k, filter=filters)
        
        # Format results
        retrieval_results = []
        for idx, score, metadata in results:
            result = RetrievalResult(
                document_id=metadata.get('doc_id', ''),
                chunk_id=metadata.get('chunk_id', ''),
                content=metadata.get('content', ''),
                score=float(score),
                metadata=metadata
            )
            retrieval_results.append(result)
            
        # Update cache
        self.cache[cache_key] = retrieval_results
        if len(self.cache) > self.cache_size:
            # Remove oldest entry
            self.cache.pop(next(iter(self.cache)))
            
        self.retrieval_complete.emit(retrieval_results)
        return retrieval_results
        
    def generate(self, query: str, context: List[str], max_length: int = 512) -> str:
        """Generate response using retrieved context"""
        # Format input
        context_str = " ".join(context)
        input_text = f"question: {query} context: {context_str}"
        
        # Tokenize
        inputs = self.tokenizer(input_text, return_tensors="pt", truncation=True, 
                              max_length=1024).to(self.device)
        
        # Generate
        with torch.no_grad():
            outputs = self.model.generate(
                **inputs,
                max_length=max_length,
                num_beams=4,
                temperature=0.7,
                do_sample=True,
                top_p=0.9
            )
            
        response = self.tokenizer.decode(outputs[0], skip_special_tokens=True)
        self.generation_complete.emit(response)
        return response
        
    def query(self, query: str, k: int = 5, filters: Dict[str, Any] = None, 
             max_length: int = 512) -> Dict[str, Any]:
        """Complete RAG query pipeline"""
        # Retrieve
        retrieval_results = self.retrieve(query, k=k, filters=filters)
        
        # Extract context
        context = [r.content for r in retrieval_results]
        
        # Generate response
        response = self.generate(query, context, max_length=max_length)
        
        # Prepare result
        result = {
            'query': query,
            'response': response,
            'sources': retrieval_results,
            'metadata': {
                'model': self.model_name,
                'k': k,
                'timestamp': datetime.now().isoformat()
            }
        }
        
        self.pipeline_complete.emit(result)
        return result
        
    def export_knowledge_base(self, path: str):
        """Export entire knowledge base"""
        self.vector_db.save(path)
        
    def import_knowledge_base(self, path: str):
        """Import knowledge base"""
        self.vector_db.load(path)


class EnterpriseRAGSystem(QObject):
    """Enterprise-grade RAG system with all features"""
    
    system_ready = pyqtSignal()
    processing_started = pyqtSignal(str)
    processing_complete = pyqtSignal(dict)
    error_occurred = pyqtSignal(str)
    
    def __init__(self, config: Dict[str, Any] = None):
        super().__init__()
        self.config = config or {}
        
        # Initialize components
        self.vector_db = VectorDatabase(
            backend=self.config.get('vector_backend', 'faiss'),
            dimension=self.config.get('embedding_dim', 768),
            use_gpu=self.config.get('use_gpu', True)
        )
        
        self.embedding_engine = EmbeddingEngine(
            model_name=self.config.get('embedding_model', 'sentence-transformers/all-mpnet-base-v2')
        )
        
        self.rag_pipeline = RAGPipeline(self.vector_db, self.embedding_engine)
        
        # Connect signals
        self.rag_pipeline.pipeline_complete.connect(self.processing_complete)
        self.rag_pipeline.error_occurred.connect(self.error_occurred)
        
        # Thread pool for parallel processing
        self.executor = ThreadPoolExecutor(max_workers=4)
        
        self.system_ready.emit()
        
    def add_document_folder(self, folder_path: str, recursive: bool = True):
        """Add all documents from folder"""
        self.processing_started.emit(f"Processing folder: {folder_path}")
        
        folder = Path(folder_path)
        if not folder.exists():
            self.error_occurred.emit(f"Folder not found: {folder_path}")
            return
            
        # Find all supported files
        patterns = ['**/*' if recursive else '*']
        files = []
        
        for pattern in patterns:
            for ext in self.rag_pipeline.doc_processor.supported_formats:
                files.extend(folder.glob(f"{pattern}{ext}"))
                
        # Process in parallel
        file_paths = [str(f) for f in files]
        self.rag_pipeline.add_documents(file_paths)
        
        self.processing_complete.emit({
            'action': 'add_folder',
            'folder': folder_path,
            'files_processed': len(files)
        })
        
    def query(self, query: str, **kwargs) -> Dict[str, Any]:
        """Query the RAG system"""
        return self.rag_pipeline.query(query, **kwargs)
        
    def export_system(self, path: str):
        """Export entire RAG system"""
        export_path = Path(path)
        export_path.mkdir(exist_ok=True)
        
        # Export vector database
        self.vector_db.save(str(export_path / "vector_db"))
        
        # Export configuration
        with open(export_path / "config.json", 'w') as f:
            json.dump(self.config, f, indent=2)
            
        # Export metadata
        metadata = {
            'export_date': datetime.now().isoformat(),
            'version': '1.0.0',
            'total_documents': self.vector_db.index.ntotal if hasattr(self.vector_db.index, 'ntotal') else 0
        }
        with open(export_path / "metadata.json", 'w') as f:
            json.dump(metadata, f, indent=2)
            
    def import_system(self, path: str):
        """Import RAG system"""
        import_path = Path(path)
        
        # Import vector database
        self.vector_db.load(str(import_path / "vector_db"))
        
        # Import configuration
        with open(import_path / "config.json", 'r') as f:
            self.config = json.load(f)


# Create global instance
ENTERPRISE_RAG = None

def initialize_rag_system(config: Dict[str, Any] = None):
    """Initialize the enterprise RAG system"""
    global ENTERPRISE_RAG
    ENTERPRISE_RAG = EnterpriseRAGSystem(config)
    return ENTERPRISE_RAG